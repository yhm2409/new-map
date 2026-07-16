# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x13, 0x4E, 0x4A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
LIGHT_TEAL_LINE = RGBColor(0xB8, 0xD3, 0xD0)
FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(blank)


def set_bg_white(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, size, color=TEAL,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return box


def add_line(slide, left, top, width, color=LIGHT_TEAL_LINE, weight=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_rect(slide, left, top, width, height, fill_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    if fill_color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp


def add_page_footer(slide, num):
    add_textbox(slide, SW - Inches(1.2), SH - Inches(0.6), Inches(0.9), Inches(0.4),
                str(num), 11, color=ORANGE, align=PP_ALIGN.RIGHT)


def add_header(slide, kicker, title):
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4), kicker, 13, color=ORANGE, bold=True)
    add_textbox(slide, Inches(0.68), Inches(0.85), Inches(11.5), Inches(0.9), title, 30, color=TEAL, bold=True)
    add_line(slide, Inches(0.7), Inches(1.7), Inches(11.9), color=LIGHT_TEAL_LINE, weight=1.25)


def content_slide(kicker, title, num):
    s = new_slide()
    set_bg_white(s)
    add_header(s, kicker, title)
    add_page_footer(s, num)
    return s


def section_slide(no, title, num):
    s = new_slide()
    set_bg_white(s)
    add_textbox(s, Inches(0.9), Inches(3.2), Inches(2), Inches(0.6), no, 20, color=ORANGE, bold=True)
    add_textbox(s, Inches(0.85), Inches(3.6), Inches(11.5), Inches(1), title, 38, color=TEAL, bold=True)
    add_line(s, Inches(0.9), Inches(4.6), Inches(2.4), color=ORANGE, weight=1.5)
    add_page_footer(s, num)
    return s


# ========== 1. 표지 ==========
s = new_slide()
set_bg_white(s)
add_rect(s, 0, SH - Inches(0.12), SW, Inches(0.12), fill_color=TEAL)
add_textbox(s, Inches(0.9), Inches(2.7), Inches(10), Inches(0.5), "MY SERVICE", 16, color=ORANGE, bold=True)
add_textbox(s, Inches(0.85), Inches(3.1), Inches(11), Inches(1.4), "나만의 서비스", 54, color=TEAL, bold=True)
add_line(s, Inches(0.9), Inches(4.45), Inches(3.2), color=ORANGE, weight=1.5)
add_textbox(s, Inches(0.9), Inches(4.65), Inches(8), Inches(0.5), "[이름]  ·  [발표일]", 16, color=TEAL)

# ========== 2. 목차 ==========
s = content_slide("CONTENTS", "목차", 2)
items = [
    ("01", "서비스 소개"),
    ("02", "서비스 구현 범위"),
    ("03", "서비스 UI 및 레퍼런스"),
    ("04", "백엔드 및 DB 구성 (ERD)"),
    ("05", "서비스 이용 시연"),
]
top = Inches(2.15)
for idx, (no, label) in enumerate(items):
    y = top + Inches(0.85) * idx
    add_textbox(s, Inches(0.8), y, Inches(1.2), Inches(0.6), no, 22, color=ORANGE, bold=True)
    add_textbox(s, Inches(1.9), y + Inches(0.03), Inches(9), Inches(0.55), label, 20, color=TEAL, bold=True)
    add_line(s, Inches(0.8), y + Inches(0.68), Inches(10.8), color=LIGHT_TEAL_LINE, weight=1.0)

# ========== slides 3~15: 헤더만 있는 빈 슬라이드 ==========
section_slide("01", "서비스 소개", 3)
content_slide("SERVICE INTRO", "기획 배경 및 타겟 사용자", 4)
content_slide("SERVICE INTRO", "핵심 기능", 5)
section_slide("02", "서비스 구현 범위", 6)
content_slide("SCOPE", "기능 명세", 7)
section_slide("03", "서비스 UI 및 레퍼런스", 8)
content_slide("USER FLOW", "사용자 플로우", 9)
content_slide("REFERENCE", "디자인 레퍼런스 및 스타일", 10)
section_slide("04", "백엔드 및 DB 구성", 11)
content_slide("DATABASE", "ERD 다이어그램", 12)
content_slide("TECH STACK", "기술 스택", 13)
section_slide("05", "서비스 이용 시연", 14)
content_slide("DEMO", "이용 시연 시나리오", 15)

out_path = r"C:\Users\권혁신\AppData\Local\Temp\claude\C--Users------claude-workspace\3fd79158-2d95-4bac-a345-3c376a6c16b1\scratchpad\나만의서비스_발표자료_템플릿.pptx"
prs.save(out_path)
print("saved:", out_path)
