from pptx import Presentation

def inspect():
    prs = Presentation("나만의서비스_발표자료_템플릿.pptx")
    print(f"Total slides in template: {len(prs.slides)}")
    for idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {idx + 1} ---")
        # Layout info
        print(f"Layout Name: {slide.slide_layout.name}")
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                print(f"Shape {shape_idx} text: '{shape.text}'")

if __name__ == "__main__":
    inspect()
