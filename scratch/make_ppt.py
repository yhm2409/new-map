import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

def build_presentation():
    prs = Presentation()
    
    # slide size: 16:9 와이드 스크린 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 6번 레이아웃은 완전 빈(Blank) 슬라이드입니다.
    blank_layout = prs.slide_layouts[6]
    
    # 헬퍼 함수: 흰색 배경 강제 설정
    def set_white_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 헬퍼 함수: 슬라이드 제목과 장표 번호 일관되게 렌더링
    def add_slide_header(slide, title_text, slide_number_str):
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10.0), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(25, 31, 40) # Toss dark text color
        
        # Slide number indicator (Toss style clean right-aligned)
        num_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.5), Inches(1.0), Inches(0.5))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = slide_number_str
        np.font.name = "Malgun Gothic"
        np.font.size = Pt(14)
        np.font.bold = True
        np.font.color.rgb = RGBColor(139, 149, 161) # Toss light gray text color

    # ----------------------------------------------------
    # SLIDE 1: 표지 (Title Slide)
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_white_bg(s1)
    
    # Logo mark
    logo_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.0), Inches(0.8))
    ltf = logo_box.text_frame
    lp = ltf.paragraphs[0]
    lp.text = "L'AROME Fine Dining"
    lp.font.name = "Malgun Gothic"
    lp.font.size = Pt(24)
    lp.font.bold = True
    lp.font.color.rgb = RGBColor(49, 130, 246) # Toss Blue
    
    # Title / Subtitle box
    title_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.3), Inches(3.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p_title = tf1.paragraphs[0]
    p_title.text = "라롬 (L'AROME) 식당 예약 및 AI 상담 플랫폼"
    p_title.font.name = "Malgun Gothic"
    p_title.font.size = Pt(38)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(25, 31, 40)
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "\n식당 메뉴와 시간을 사전에 조율하여 예약하는 스마트 모바일 웹앱 서비스"
    p_sub.font.name = "Malgun Gothic"
    p_sub.font.size = Pt(18)
    p_sub.font.bold = True
    p_sub.font.color.rgb = RGBColor(78, 89, 104)

    # ----------------------------------------------------
    # SLIDE 2: 목차 (Table of Contents)
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_white_bg(s2)
    add_slide_header(s2, "목차 (Table of Contents)", "Index")
    
    txBoxIndex = s2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
    tfi = txBoxIndex.text_frame
    tfi.word_wrap = True
    
    index_items = [
      "1. 서비스 소개 - 식당 사전 예약 플랫폼의 목적 및 필요성",
      "2. 서비스 구현 범위 - Client, Server, AI 가용 범위 개괄",
      "3. 서비스 UI 및 레퍼런스 - TOSS 디자인 상속 및 작업 프로세스",
      "4. 백엔드 구성 - Next.js Route Handlers 및 보안 인증 관리 API",
      "5. DB 구성 (ERD) - Supabase PostgreSQL 관계형 스키마 설계",
      "6. 서비스 이용 시연 - 실시간 2D 테이블 선점 및 중복 예약 방지",
      "7. 서비스 이용 시연 - AI 레스토랑 컨시어지 챗봇 자동화 상담"
    ]
    
    for idx, item in enumerate(index_items):
        p = tfi.paragraphs[0] if idx == 0 else tfi.add_paragraph()
        p.text = item
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(25, 31, 40)
        p.space_after = Pt(12)

    # ----------------------------------------------------
    # SLIDE 3: 1. 서비스 소개 (Service Introduction)
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_white_bg(s3)
    add_slide_header(s3, "1. 서비스 소개", "01")
    
    txBox3 = s3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    
    p_intro_1 = tf3.paragraphs[0]
    p_intro_1.text = "• 식당 메뉴와 시간을 미리 예약하고 조율하는 모바일 플랫폼"
    p_intro_1.font.name = "Malgun Gothic"
    p_intro_1.font.size = Pt(20)
    p_intro_1.font.bold = True
    p_intro_1.font.color.rgb = RGBColor(25, 31, 40)
    
    p_intro_1_sub = tf3.add_paragraph()
    p_intro_1_sub.text = "   - 바쁜 현대인들의 방문 대기 시간을 줄이고 노쇼(No-Show)를 방지하는 실질적 가치 창출.\n   - 소비자는 사전 예약을 통해 기다림 없는 미식을 즐기고, 가맹점은 효율적으로 테이블 회전율을 조율할 수 있습니다."
    p_intro_1_sub.font.name = "Malgun Gothic"
    p_intro_1_sub.font.size = Pt(15)
    p_intro_1_sub.font.color.rgb = RGBColor(78, 89, 104)
    p_intro_1_sub.space_after = Pt(20)
    
    p_intro_2 = tf3.add_paragraph()
    p_intro_2.text = "• 메뉴 및 일정 사전 선택의 이점"
    p_intro_2.font.name = "Malgun Gothic"
    p_intro_2.font.size = Pt(20)
    p_intro_2.font.bold = True
    p_intro_2.font.color.rgb = RGBColor(25, 31, 40)
    
    p_intro_2_sub = tf3.add_paragraph()
    p_intro_2_sub.text = "   - 방문 날짜, 디너 타임, 상세 코스 요리까지 웹앱을 통해 미리 조율 완료.\n   - 매장 입장 후 추가 주문이나 대기 없이 사전에 확정된 파인다이닝 코스가 자연스럽게 서빙되는 프리미엄 동선 제공."
    p_intro_2_sub.font.name = "Malgun Gothic"
    p_intro_2_sub.font.size = Pt(15)
    p_intro_2_sub.font.color.rgb = RGBColor(78, 89, 104)

    # ----------------------------------------------------
    # SLIDE 4: 2. 서비스 구현 범위 (Implementation Scope)
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_white_bg(s4)
    add_slide_header(s4, "2. 서비스 구현 범위", "02")
    
    txBox4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    
    p_scope_1 = tf4.paragraphs[0]
    p_scope_1.text = "• 클라이언트 환경 (Client Scope)"
    p_scope_1.font.name = "Malgun Gothic"
    p_scope_1.font.size = Pt(18)
    p_scope_1.font.bold = True
    p_scope_1.font.color.rgb = RGBColor(25, 31, 40)
    
    p_scope_1_sub = tf4.add_paragraph()
    p_scope_1_sub.text = "   - **모바일 앱 뷰포트**: 480px 섀도우 기기 프레임 적용으로 스마트폰 접속 시 화면에 최적 핏팅되는 레이아웃.\n   - **예약 3단계 흐름**: 일정/인원 선택 ➡️ 2D 테이블 좌석 지정 ➡️ 코스 요리 메뉴 선택(장바구니) 및 모의 결제창 연동.\n   - **Sticky 하단 네비게이션**: 예약 홈, 내 예약 현황, AI 챗봇 3대 탭 동적 스위칭 구현."
    p_scope_1_sub.font.name = "Malgun Gothic"
    p_scope_1_sub.font.size = Pt(14)
    p_scope_1_sub.font.color.rgb = RGBColor(78, 89, 104)
    p_scope_1_sub.space_after = Pt(15)
    
    p_scope_2 = tf4.add_paragraph()
    p_scope_2.text = "• 백엔드 및 관리자 대시보드 (Server & Admin Scope)"
    p_scope_2.font.name = "Malgun Gothic"
    p_scope_2.font.size = Pt(18)
    p_scope_2.font.bold = True
    p_scope_2.font.color.rgb = RGBColor(25, 31, 40)
    
    p_scope_2_sub = tf4.add_paragraph()
    p_scope_2_sub.text = "   - **Next.js API Routes**: 인증(Auth), 예약 및 주문 내역, 좌석 실시간 예약 상태 제공 백엔드 핸들러.\n   - **사장님 대시보드 (어드민)**: 데스크톱 가로 모드 해상도 반응형 제공. 예약 건 승인, 반려 및 Cascade 물리 영구 삭제 처리."
    p_scope_2_sub.font.name = "Malgun Gothic"
    p_scope_2_sub.font.size = Pt(14)
    p_scope_2_sub.font.color.rgb = RGBColor(78, 89, 104)

    # ----------------------------------------------------
    # SLIDE 5: 3. 서비스 UI 및 레퍼런스 (UI & Process)
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_white_bg(s5)
    add_slide_header(s5, "3. 서비스 UI 및 레퍼런스 (작업 프로세스)", "03")
    
    # Text info on left
    txBox5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    tf5 = txBox5.text_frame
    tf5.word_wrap = True
    
    p_ui_1 = tf5.paragraphs[0]
    p_ui_1.text = "• 토스(TOSS) 디자인 레퍼런스 지향"
    p_ui_1.font.name = "Malgun Gothic"
    p_ui_1.font.size = Pt(18)
    p_ui_1.font.bold = True
    p_ui_1.font.color.rgb = RGBColor(25, 31, 40)
    
    p_ui_1_sub = tf5.add_paragraph()
    p_ui_1_sub.text = "   - 명료한 블루(#3182f6)를 메인 색상으로 차용해 깔끔하고 웅장한 가독성 확보.\n   - 플랫한 둥근 카드 레이아웃과 맑은 디자인 컴포넌트 결합.\n   - 높이가 제한된 모바일 앱 화면 내에서 메뉴판 스크롤과 카드를 상하 2단 배치해 공간 효율을 최대화했습니다."
    p_ui_1_sub.font.name = "Malgun Gothic"
    p_ui_1_sub.font.size = Pt(13)
    p_ui_1_sub.font.color.rgb = RGBColor(78, 89, 104)
    p_ui_1_sub.space_after = Pt(15)
    
    p_ui_2 = tf5.add_paragraph()
    p_ui_2.text = "• Iterative (점진적) 개발 프로세스"
    p_ui_2.font.name = "Malgun Gothic"
    p_ui_2.font.size = Pt(18)
    p_ui_2.font.bold = True
    p_ui_2.font.color.rgb = RGBColor(25, 31, 40)
    
    p_ui_2_sub = tf5.add_paragraph()
    p_ui_2_sub.text = "   - 1단계: DB 스키마 설계 및 Prisma 배포 ➡️ 2단계: 모바일 예약 동선 구현 ➡️ 3단계: AI 챗봇 탑재 ➡️ 4단계: 피드백 반영 롤백.\n   - 영화 예매 도메인 전환 피드백을 접수하고, 식당 예약 뼈대로 신속히 원상 복구 및 동기화하는 기민한(Agile) 대처 프로세스를 수행했습니다."
    p_ui_2_sub.font.name = "Malgun Gothic"
    p_ui_2_sub.font.size = Pt(13)
    p_ui_2_sub.font.color.rgb = RGBColor(78, 89, 104)

    # Insert WebP reference image converting to PNG on the fly
    try:
        toss_webp = r"c:\Users\yhm2409\Desktop\map\toss\imgi_61_73ab268a060611ed89c61329738d4db8_small.webp"
        temp_png = r"c:\Users\yhm2409\Desktop\map\toss\temp_toss_ref2.png"
        if os.path.exists(toss_webp):
            with Image.open(toss_webp) as img:
                img.save(temp_png, "PNG")
            s5.shapes.add_picture(temp_png, Inches(8.5), Inches(1.8), width=Inches(3.8))
    except Exception as e:
        print("Toss WebP convert failed")

    # ----------------------------------------------------
    # SLIDE 6: 4. 백엔드 구성 (Backend API Architecture)
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_white_bg(s6)
    add_slide_header(s6, "4. 백엔드 구성 (Backend Architecture)", "04")
    
    txBox6 = s6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    tf6 = txBox6.text_frame
    tf6.word_wrap = True
    
    p_be_1 = tf6.paragraphs[0]
    p_be_1.text = "• Next.js App Router API Routes 설계"
    p_be_1.font.name = "Malgun Gothic"
    p_be_1.font.size = Pt(20)
    p_be_1.font.bold = True
    p_be_1.font.color.rgb = RGBColor(25, 31, 40)
    
    p_be_1_sub = tf6.add_paragraph()
    p_be_1_sub.text = "   - **`/api/auth/login` & `/api/auth/signup`**: 암호화 및 유저 인증 관리 엔드포인트.\n   - **`/api/reservations`**: POST 호출 시 DB 트랜잭션을 작동하여 좌석의 동시 예약을 원천 제어.\n   - **`/api/reservations/my` & `/api/reservations/admin`**: 유저 고유 ID 기준 예약 조회 및 전체 대시보드 리스트 출력 API.\n   - **`/api/reservations/[id]`**: PATCH(상태변경: 승인/거절) 및 DELETE(예약 물리 완전 Cascade 삭제) 처리."
    p_be_1_sub.font.name = "Malgun Gothic"
    p_be_1_sub.font.size = Pt(14)
    p_be_1_sub.font.color.rgb = RGBColor(78, 89, 104)
    p_be_1_sub.space_after = Pt(20)
    
    p_be_2 = tf6.add_paragraph()
    p_be_2.text = "• AI 추론 연동 라우트 (`/api/chat`)"
    p_be_2.font.name = "Malgun Gothic"
    p_be_2.font.size = Pt(20)
    p_be_2.font.bold = True
    p_be_2.font.color.rgb = RGBColor(25, 31, 40)
    
    p_be_2_sub = tf6.add_paragraph()
    p_be_2_sub.text = "   - 클라이언트에서 전송된 채팅 히스토리를 받아 서버사이드에서 GROQ API에 위임 처리.\n   - API Key 노출을 원천 방지하기 위해 클라이언트가 아닌 Next.js Server Side에서 안전하게 API Key와 System Prompt를 병합하여 통신."
    p_be_2_sub.font.name = "Malgun Gothic"
    p_be_2_sub.font.size = Pt(14)
    p_be_2_sub.font.color.rgb = RGBColor(78, 89, 104)

    # ----------------------------------------------------
    # SLIDE 7: 5. DB 구성 (ERD Schema Diagram)
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_white_bg(s7)
    add_slide_header(s7, "5. DB 구성 (ERD)", "05")
    
    # Text info on left
    txBox7 = s7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
    tf7 = txBox7.text_frame
    tf7.word_wrap = True
    
    p_db_1 = tf7.paragraphs[0]
    p_db_1.text = "• Prisma ORM 5 & Supabase DB 결합"
    p_db_1.font.name = "Malgun Gothic"
    p_db_1.font.size = Pt(18)
    p_db_1.font.bold = True
    p_db_1.font.color.rgb = RGBColor(25, 31, 40)
    
    p_db_1_sub = tf7.add_paragraph()
    p_db_1_sub.text = "   - Supabase PostgreSQL 클라우드 원격 연결 구조로 데이터 신뢰성 확보.\n   - 특수 문자(#, %) URL 인코딩 처리를 거친 DATABASE_URL 환경 변수로 연결 안정화."
    p_db_1_sub.font.name = "Malgun Gothic"
    p_db_1_sub.font.size = Pt(13)
    p_db_1_sub.font.color.rgb = RGBColor(78, 89, 104)
    p_db_1_sub.space_after = Pt(15)
    
    p_db_2 = tf7.add_paragraph()
    p_db_2.text = "• 테이블 설계 및 Cascade 제약"
    p_db_2.font.name = "Malgun Gothic"
    p_db_2.font.size = Pt(18)
    p_db_2.font.bold = True
    p_db_2.font.color.rgb = RGBColor(25, 31, 40)
    
    p_db_2_sub = tf7.add_paragraph()
    p_db_2_sub.text = "   - **User 테이블**: id(PK), loginId, name, password, role(관리자/사용자 구분) 관리.\n   - **Reservation 테이블**: id(PK), userId(FK), 날짜, 시간, 지정된 테이블 식별자(tableId) 보관.\n   - **ReservationItem 테이블**: id(PK), reservationId(FK), 요리 품목 가격 및 주문 수량 관리.\n   - 어드민 단에서 예약을 삭제(DELETE)하면 연결된 ReservationItem 레코드가 Cascade Delete 룰에 의해 안전하게 연쇄 자동 소멸."
    p_db_2_sub.font.name = "Malgun Gothic"
    p_db_2_sub.font.size = Pt(13)
    p_db_2_sub.font.color.rgb = RGBColor(78, 89, 104)

    # Insert user provided ERD image (Updated with the newest media file)
    erd_artifact_path = r"C:\Users\yhm2409\.gemini\antigravity\brain\514b8155-baf4-4d6f-923a-5def5c7653b5\media__1784168928988.png"
    if os.path.exists(erd_artifact_path):
        s7.shapes.add_picture(erd_artifact_path, Inches(6.8), Inches(1.8), width=Inches(5.7))

    # ----------------------------------------------------
    # SLIDE 8: 6. 서비스 이용 시연 (실시간 중복 예약 방지)
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_white_bg(s8)
    add_slide_header(s8, "6. 서비스 이용 시연 (실시간 중복 예약 방지)", "06")
    
    txBox8 = s8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    tf8 = txBox8.text_frame
    tf8.word_wrap = True
    
    p_demo_1 = tf8.paragraphs[0]
    p_demo_1.text = "• 날짜 및 시간대별 실시간 좌석 점유 상태 조회"
    p_demo_1.font.name = "Malgun Gothic"
    p_demo_1.font.size = Pt(20)
    p_demo_1.font.bold = True
    p_demo_1.font.color.rgb = RGBColor(25, 31, 40)
    
    p_demo_1_sub = tf8.add_paragraph()
    p_demo_1_sub.text = "   - 예약 일정 단계에서 날짜와 디너 시간대를 클릭해 변경하는 즉시 `fetchReservedTables` API가 작동.\n   - 해당 슬롯에 이미 완료된 예약이 존재하는 좌석 테이블 ID 목록을 데이터베이스로부터 정확하게 호출합니다."
    p_demo_1_sub.font.name = "Malgun Gothic"
    p_demo_1_sub.font.size = Pt(15)
    p_demo_1_sub.font.color.rgb = RGBColor(78, 89, 104)
    p_demo_1_sub.space_after = Pt(20)
    
    p_demo_2 = tf8.add_paragraph()
    p_demo_2.text = "• SeatMap 좌석 비활성화 및 동시 예약 차단"
    p_demo_2.font.name = "Malgun Gothic"
    p_demo_2.font.size = Pt(20)
    p_demo_2.font.bold = True
    p_demo_2.font.color.rgb = RGBColor(25, 31, 40)
    
    p_demo_2_sub = tf8.add_paragraph()
    p_demo_2_sub.text = "   - 이미 선점된 좌석 테이블(Table 1~6)은 SeatMap에서 연회색 비활성화 상태로 락 처리되어 마우스/터치 선택이 불가능.\n   - DB 트랜잭션 수준에서 동시 접수 충돌을 완벽 제어하여 1테이블당 중복 예약을 원천적으로 방지합니다."
    p_demo_2_sub.font.name = "Malgun Gothic"
    p_demo_2_sub.font.size = Pt(15)
    p_demo_2_sub.font.color.rgb = RGBColor(78, 89, 104)

    # ----------------------------------------------------
    # SLIDE 9: 6. 서비스 이용 시연 (AI 챗봇 및 어드민 연계)
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_white_bg(s9)
    add_slide_header(s9, "6. 서비스 이용 시연 (AI 레스토랑 챗봇)", "07")
    
    txBox9 = s9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    tf9 = txBox9.text_frame
    tf9.word_wrap = True
    
    p_chat_1 = tf9.paragraphs[0]
    p_chat_1.text = "• GROQ Llama-3.3 기반 식당 정보 100% 동기화 상담"
    p_chat_1.font.name = "Malgun Gothic"
    p_chat_1.font.size = Pt(20)
    p_chat_1.font.bold = True
    p_chat_1.font.color.rgb = RGBColor(25, 31, 40)
    
    p_chat_1_sub = tf9.add_paragraph()
    p_chat_1_sub.text = "   - 예매 웹 폼 상에 실제 표시되는 요리 정보(아란치니, 한우 스테이크 등) 및 6개 테이블의 특징(창가석, 룸 등), 무료 발레파킹 혜택을 챗봇의 지식 컨텍스트에 고스란히 이식.\n   - 소비자가 창가석이나 주차 요금, 메뉴 종류에 대해 물어보면 매장의 실제 세부 사양과 완벽히 동기화된 일관성 있는 답변을 산출합니다."
    p_chat_1_sub.font.name = "Malgun Gothic"
    p_chat_1_sub.font.size = Pt(15)
    p_chat_1_sub.font.color.rgb = RGBColor(78, 89, 104)
    p_chat_1_sub.space_after = Pt(20)
    
    p_chat_2 = tf9.add_paragraph()
    p_chat_2.text = "• 어드민 예약 처리 연동 실시간 동기화"
    p_chat_2.font.name = "Malgun Gothic"
    p_chat_2.font.size = Pt(20)
    p_chat_2.font.bold = True
    p_chat_2.font.color.rgb = RGBColor(25, 31, 40)
    
    p_chat_2_sub = tf9.add_paragraph()
    p_chat_2_sub.text = "   - 사장님이 어드민 모드에서 승인 완료된 예약을 영구 삭제 시, 데이터베이스 물리 제거와 동시에 고객의 내 예약 확인 탭에서도 실시간 삭제 처리 연계.\n   - 삭제와 동시에 해당 좌석 락이 즉각 해제되어 SeatMap 상에서 실시간으로 다시 예약 가능 상태로 복구됩니다."
    p_chat_2_sub.font.name = "Malgun Gothic"
    p_chat_2_sub.font.size = Pt(15)
    p_chat_2_sub.font.color.rgb = RGBColor(78, 89, 104)

    # ----------------------------------------------------
    # SLIDE 10: Q&A / 감사합니다 (Conclusion Slide)
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_white_bg(s10)
    
    txBox10 = s10.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(3.0))
    tf10 = txBox10.text_frame
    tf10.word_wrap = True
    
    p_end = tf10.paragraphs[0]
    p_end.text = "경청해 주셔서 감사합니다"
    p_end.font.name = "Malgun Gothic"
    p_end.font.size = Pt(44)
    p_end.font.bold = True
    p_end.font.color.rgb = RGBColor(49, 130, 246) # Toss Blue
    
    p_end_sub = tf10.add_paragraph()
    p_end_sub.text = "\n라롬 (L'AROME) 식당 예약 및 AI 상담 플랫폼 Q&A"
    p_end_sub.font.name = "Malgun Gothic"
    p_end_sub.font.size = Pt(20)
    p_end_sub.font.bold = True
    p_end_sub.font.color.rgb = RGBColor(78, 89, 104)

    # Save to user desktop directory (Changed to v2 to prevent PPT program lock collision)
    dest_path = r"c:\Users\yhm2409\Desktop\map\L_AROME_Presentation_v2.pptx"
    prs.save(dest_path)
    
    # clean up temp image
    if os.path.exists(temp_png):
        os.remove(temp_png)

if __name__ == "__main__":
    build_presentation()
