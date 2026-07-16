# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from PIL import Image

# 브랜드 컬러 코드 정의
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BRAND_BLUE = RGBColor(0x1C, 0x60, 0xEF)     # #1C60EF
SOFT_BLUE = RGBColor(0xDB, 0xE8, 0xFE)      # #DBE8FE
DARK_TEXT = RGBColor(0x19, 0x1F, 0x28)      # #191F28 (Toss Dark)
GRAY_TEXT = RGBColor(0x4E, 0x59, 0x68)      # #4E5968 (Toss Gray)
LIGHT_LINE = RGBColor(0xE5, 0xE8, 0xEB)      # #E5E8EB (구분선)

# 지정된 폰트 정의
FONT_TITLE = "KIMM_오피스체_Bold"
FONT_BODY = "에이투지체"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank_layout = prs.slide_layouts[6]

def new_slide():
    s = prs.slides.add_slide(blank_layout)
    # 흰색 배경 강제 설정
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s

# ----------------------------------------------------
# 셰이프 및 텍스트 추가 헬퍼 함수들 (도형 위주 디자인)
# ----------------------------------------------------

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    return shp

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, left, top, width, height, text, size, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font=FONT_BODY, vertical=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box

def add_header(slide, kicker, title):
    # Kicker (본문 폰트, 브랜드 블루 색상)
    add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.4), kicker, 13, color=BRAND_BLUE, bold=True, font=FONT_BODY)
    # Title (KIMM 타이틀 폰트, 어두운 색상)
    add_text(slide, Inches(0.78), Inches(0.82), Inches(11.5), Inches(0.9), title, 32, color=DARK_TEXT, bold=True, font=FONT_TITLE)
    # 구분 가로선
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.8), Inches(1.65), Inches(12.5), Inches(1.65))
    ln.line.color.rgb = LIGHT_LINE
    ln.line.width = Pt(1.5)

def add_footer(slide, page_num):
    add_text(slide, SW - Inches(1.5), SH - Inches(0.6), Inches(0.9), Inches(0.4),
             str(page_num), 11, color=GRAY_TEXT, align=PP_ALIGN.RIGHT, font=FONT_BODY)

def create_section_slide(number_str, title, page_num):
    s = new_slide()
    # 왼쪽 브랜드 블루 디자인 띠 세로로 렌더링
    add_rect(s, 0, 0, Inches(0.3), SH, BRAND_BLUE)
    # 웅장한 번호
    add_text(s, Inches(1.2), Inches(2.8), Inches(3), Inches(0.6), number_str, 24, color=BRAND_BLUE, bold=True, font=FONT_BODY)
    # 대제목
    add_text(s, Inches(1.15), Inches(3.3), Inches(10), Inches(1.2), title, 42, color=DARK_TEXT, bold=True, font=FONT_TITLE)
    # 하단 얇은 블루 라인
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), Inches(4.5), Inches(4.2), Inches(4.5))
    ln.line.color.rgb = BRAND_BLUE
    ln.line.width = Pt(2.0)
    add_footer(s, page_num)
    return s

def build():
    # ----------------------------------------------------
    # SLIDE 1: 표지 (Title Slide)
    # ----------------------------------------------------
    s1 = new_slide()
    # 중앙 브랜드 블루 대형 카드 배치
    add_rounded_rect(s1, Inches(0.8), Inches(1.5), Inches(11.73), Inches(4.2), fill_color=BRAND_BLUE)
    
    # 카드 내 텍스트
    add_text(s1, Inches(1.4), Inches(2.2), Inches(10), Inches(0.5), "L'AROME FINE DINING", 16, color=SOFT_BLUE, bold=True, font=FONT_BODY)
    add_text(s1, Inches(1.35), Inches(2.7), Inches(10.5), Inches(1.8), "라롬 식당 예약 및 AI 상담 플랫폼", 46, color=WHITE, bold=True, font=FONT_TITLE)
    
    # 카드 밑줄 데코
    ln1 = s1.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.4), Inches(4.1), Inches(4.8), Inches(4.1))
    ln1.line.color.rgb = WHITE
    ln1.line.width = Pt(2.0)
    
    # 하단 메타 정보
    add_text(s1, Inches(1.4), Inches(4.4), Inches(10), Inches(0.5), "식당 메뉴와 시간을 사전에 조율하는 예약 웹앱", 16, color=SOFT_BLUE, font=FONT_BODY)
    add_text(s1, Inches(1.4), Inches(4.8), Inches(10), Inches(0.5), "발표자: [이름]  ·  개발 프로젝트 완료 보고", 14, color=WHITE, font=FONT_BODY)
    add_footer(s1, 1)

    # ----------------------------------------------------
    # SLIDE 2: CONTENTS (목차)
    # ----------------------------------------------------
    s2 = new_slide()
    add_header(s2, "CONTENTS", "목차")
    
    items = [
        ("01", "서비스 소개", "식당 예약 플랫폼의 목적 및 필요성"),
        ("02", "서비스 구현 범위", "Client, Server, AI 가용 범위 개괄"),
        ("03", "서비스 UI 및 레퍼런스", "TOSS 디자인 스타일 및 작업 프로세스"),
        ("04", "백엔드 및 DB 구성", "Next.js Route Handlers 및 DB ERD"),
        ("05", "서비스 이용 시연", "중복 예약 방지 및 AI 실시간 상담")
    ]
    
    # 5개의 깔끔한 가로형 카드 박스 배치 (도식화)
    card_width = Inches(11.73)
    card_height = Inches(0.85)
    top_base = Inches(2.0)
    
    for idx, (num, title, desc) in enumerate(items):
        y_pos = top_base + Inches(0.95) * idx
        # 연블루 배경 라운드 박스
        add_rounded_rect(s2, Inches(0.8), y_pos, card_width, card_height, SOFT_BLUE)
        
        # 순번
        add_text(s2, Inches(1.2), y_pos + Inches(0.2), Inches(0.8), Inches(0.5), num, 20, color=BRAND_BLUE, bold=True, font=FONT_BODY)
        # 제목
        add_text(s2, Inches(2.0), y_pos + Inches(0.2), Inches(4.0), Inches(0.5), title, 18, color=DARK_TEXT, bold=True, font=FONT_TITLE)
        # 설명
        add_text(s2, Inches(6.0), y_pos + Inches(0.25), Inches(6.0), Inches(0.5), desc, 14, color=GRAY_TEXT, font=FONT_BODY)

    add_footer(s2, 2)

    # ----------------------------------------------------
    # SLIDE 3: 01. 서비스 소개 간지
    # ----------------------------------------------------
    create_section_slide("01", "서비스 소개", 3)

    # ----------------------------------------------------
    # SLIDE 4: 서비스 소개 (기획 배경) - 세로 카드 3개 도식화
    # ----------------------------------------------------
    s4 = new_slide()
    add_header(s4, "SERVICE INTRO", "기획 배경 및 해결하고자 하는 문제")
    
    intro_cards = [
        ("• 대기 피로도 극복", "식당 입장을 위해 장시간 현장에서 대기하는 비효율성을 극복하여 미식 소비의 만족도를 높입니다."),
        ("• 예약 부도(No-Show) 방지", "방문 일정뿐만 아니라 세부 식사 메뉴까지 모바일을 통해 선주문함으로써 매장의 노쇼 리스크를 차단합니다."),
        ("• 정보 불일치 해소", "창가석, 프라이빗 룸 등 고객이 실제 선호하는 좌석 형태를 2D 배치도로 확인하여 정보 격차를 완전히 해소합니다.")
    ]
    
    card_w = Inches(3.7)
    card_h = Inches(4.2)
    left_start = Inches(0.8)
    gap = Inches(0.3)
    
    for idx, (title, content) in enumerate(intro_cards):
        x_pos = left_start + (card_w + gap) * idx
        # 연블루 박스
        add_rounded_rect(s4, x_pos, Inches(2.2), card_w, card_h, SOFT_BLUE)
        
        # 데코 탑 라인
        add_rect(s4, x_pos, Inches(2.2), card_w, Inches(0.12), BRAND_BLUE)
        
        # 텍스트 기입
        add_text(s4, x_pos + Inches(0.3), Inches(2.6), card_w - Inches(0.6), Inches(0.6), title, 18, color=BRAND_BLUE, bold=True, font=FONT_TITLE)
        add_text(s4, x_pos + Inches(0.3), Inches(3.3), card_w - Inches(0.6), Inches(2.8), content, 14, color=DARK_TEXT, font=FONT_BODY)

    add_footer(s4, 4)

    # ----------------------------------------------------
    # SLIDE 5: 서비스 소개 (서비스 가치) - 3단계 흐름도 도식화
    # ----------------------------------------------------
    s5 = new_slide()
    add_header(s5, "SERVICE VALUE", "서비스 가치 및 제공 기능 개요")
    
    steps = [
        ("STEP 01", "방문일 및 일정 지정", "캘린더 UI를 통해 원하는 예약 날짜와 인원, 디너 시간대를 정밀 세팅합니다."),
        ("STEP 02", "가상 테이블 지정", "실시간 점유 상태를 반영하여 창가석, 룸 등 원하는 테이블을 SeatMap 지도에서 선택합니다."),
        ("STEP 03", "코스 메뉴 담기 및 결제", "매칭된 파인다이닝 코스 요리를 장바구니에 사전에 담고 모의 선결제를 확정합니다.")
    ]
    
    step_w = Inches(3.6)
    step_h = Inches(3.8)
    s_left = Inches(0.8)
    s_gap = Inches(0.4)
    
    for idx, (step_num, title, desc) in enumerate(steps):
        x_pos = s_left + (step_w + s_gap) * idx
        # 라운드 박스
        add_rounded_rect(s5, x_pos, Inches(2.4), step_w, step_h, SOFT_BLUE)
        
        # 스텝 숫자 (화이트 텍스트 올릴 원형/사각형 모양 데코)
        add_rect(s5, x_pos + Inches(0.3), Inches(2.7), Inches(1.1), Inches(0.4), BRAND_BLUE)
        add_text(s5, x_pos + Inches(0.3), Inches(2.75), Inches(1.1), Inches(0.4), step_num, 13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_BODY)
        
        # 텍스트
        add_text(s5, x_pos + Inches(0.3), Inches(3.4), step_w - Inches(0.6), Inches(0.6), title, 18, color=BRAND_BLUE, bold=True, font=FONT_TITLE)
        add_text(s5, x_pos + Inches(0.3), Inches(4.1), step_w - Inches(0.6), Inches(1.8), desc, 14, color=DARK_TEXT, font=FONT_BODY)
        
        # 연결 화살표 데코 (마지막 카드가 아니면 그림)
        if idx < 2:
            arrow_x = x_pos + step_w + Inches(0.05)
            # 연결 가로 실선
            ln = s5.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, arrow_x, Inches(4.3), arrow_x + Inches(0.3), Inches(4.3))
            ln.line.color.rgb = BRAND_BLUE
            ln.line.width = Pt(2.0)

    add_footer(s5, 5)

    # ----------------------------------------------------
    # SLIDE 6: 02. 서비스 구현 범위 간지
    # ----------------------------------------------------
    create_section_slide("02", "서비스 구현 범위", 6)

    # ----------------------------------------------------
    # SLIDE 7: 구현 범위 (SCOPE) - 좌우 분할 박스 도식화
    # ----------------------------------------------------
    s7 = new_slide()
    add_header(s7, "DEVELOPMENT SCOPE", "Client, Server, AI 구현 명세")
    
    # Left Card: Client
    add_rounded_rect(s7, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5), SOFT_BLUE)
    add_rect(s7, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.12), BRAND_BLUE)
    add_text(s7, Inches(1.1), Inches(2.4), Inches(5.0), Inches(0.5), "클라이언트 (Client UI/UX)", 20, color=BRAND_BLUE, bold=True, font=FONT_TITLE)
    
    client_bullets = [
        "• 모바일 반응형 뷰포트 & 480px 기기 프레임 적용",
        "• Sticky Bottom Navigation (예약 홈/내역/챗봇)",
        "• 3단계 가이드 예약 모듈 및 사전 메뉴 담기 장바구니",
        "• 결제 금액 실시간 합산 & 모의 결제 연동"
    ]
    for b_idx, bullet in enumerate(client_bullets):
        add_text(s7, Inches(1.1), Inches(3.1) + Inches(0.7) * b_idx, Inches(5.0), Inches(0.6), bullet, 14, color=DARK_TEXT, font=FONT_BODY)

    # Right Card: Server & AI
    add_rounded_rect(s7, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.5), SOFT_BLUE)
    add_rect(s7, Inches(6.9), Inches(2.0), Inches(5.6), Inches(0.12), BRAND_BLUE)
    add_text(s7, Inches(7.2), Inches(2.4), Inches(5.0), Inches(0.5), "백엔드 및 AI (Server & AI)", 20, color=BRAND_BLUE, bold=True, font=FONT_TITLE)
    
    server_bullets = [
        "• Next.js API Routes (로그인, 예약 상태, 챗봇)",
        "• 실시간 중복 예약을 방지하는 날짜/시간별 좌석 점유 조회 API",
        "• 사장님 전용 대시보드 (예약 승인, 거절 및 Cascade 완전삭제)",
        "• GROQ SDK Llama-3.3 LLM 연계 실시간 자동 상담 챗봇"
    ]
    for b_idx, bullet in enumerate(server_bullets):
        add_text(s7, Inches(7.2), Inches(3.1) + Inches(0.7) * b_idx, Inches(5.0), Inches(0.6), bullet, 14, color=DARK_TEXT, font=FONT_BODY)

    add_footer(s7, 7)

    # ----------------------------------------------------
    # SLIDE 8: 03. 서비스 UI 및 레퍼런스 간지
    # ----------------------------------------------------
    create_section_slide("03", "서비스 UI 및 레퍼런스", 8)

    # ----------------------------------------------------
    # SLIDE 9: UI 및 레퍼런스 (REFERENCE)
    # ----------------------------------------------------
    s9 = new_slide()
    add_header(s9, "UI & REFERENCE", "TOSS 디자인 스타일 및 애자일 작업 프로세스")
    
    # Left Information card
    add_rounded_rect(s9, Inches(0.8), Inches(2.0), Inches(7.3), Inches(4.5), SOFT_BLUE)
    add_rect(s9, Inches(0.8), Inches(2.0), Inches(7.3), Inches(0.12), BRAND_BLUE)
    
    add_text(s9, Inches(1.1), Inches(2.4), Inches(6.7), Inches(0.5), "토스(TOSS) 디자인 레퍼런스 적용", 18, color=BRAND_BLUE, bold=True, font=FONT_TITLE)
    
    ref_bullets = [
        "• 깔끔하고 신뢰감을 주는 Toss Blue (#1C60EF) 브랜드 테마 적용.",
        "• 테두리를 최소화하고 둥글고 정갈한 플랫 카드 레이아웃 적용.",
        "• 모바일 하단 네비게이션 가이드를 완비하여 엄지 터치 안정성 확보.",
        "• 애자일 점진적 작업 프로세스를 거쳐 뼈대를 완성하고, 기획 피드백에 맞춰 식당 예약 지식 베이스 및 UI 통합을 영구 정밀화했습니다."
    ]
    for idx, bullet in enumerate(ref_bullets):
        add_text(s9, Inches(1.1), Inches(3.0) + Inches(0.8) * idx, Inches(6.7), Inches(0.75), bullet, 14, color=DARK_TEXT, font=FONT_BODY)

    # Right Image insertion
    try:
        toss_webp = r"c:\Users\yhm2409\Desktop\map\toss\imgi_61_73ab268a060611ed89c61329738d4db8_small.webp"
        temp_png = r"c:\Users\yhm2409\Desktop\map\toss\temp_toss_ref_blue.png"
        if os.path.exists(toss_webp):
            with Image.open(toss_webp) as img:
                img.save(temp_png, "PNG")
            s9.shapes.add_picture(temp_png, Inches(8.7), Inches(2.0), width=Inches(3.8))
    except Exception as e:
        print("Toss Image insertion failed:", e)

    add_footer(s9, 9)

    # ----------------------------------------------------
    # SLIDE 10: 04. 백엔드 및 DB 구성 간지
    # ----------------------------------------------------
    create_section_slide("04", "백엔드 및 DB 구성", 10)

    # ----------------------------------------------------
    # SLIDE 11: 백엔드 구성 (BACKEND) - 가로형 카드 4개 타임라인 도식화
    # ----------------------------------------------------
    s11 = new_slide()
    add_header(s11, "BACKEND ARCHITECTURE", "Next.js API Routes 및 AI 추론 연동")
    
    be_cards = [
        ("Auth API", "/api/auth/login", "사용자 계정 생성, 로그인 암호화 및 유저 권한(사용자/어드민) 세션 토큰화 관리"),
        ("Reservation API", "/api/reservations", "방문 날짜 및 시간대별 실시간 좌석 점유 조회 및 DB 점유 접수 처리"),
        ("Admin Control", "/api/reservations/[id]", "사장님 모드에서의 예약 승인/거절 상태 변경 및 완전 Cascade Delete 삭제"),
        ("AI Chat API", "/api/chat", "API Key의 브라우저 노출을 방지하고 서버사이드에서 안전하게 LLM 추론 위임")
    ]
    
    card_w = Inches(2.7)
    card_h = Inches(4.3)
    left_s = Inches(0.8)
    gap_s = Inches(0.3)
    
    for idx, (title, path, desc) in enumerate(be_cards):
        x_pos = left_s + (card_w + gap_s) * idx
        # 박스
        add_rounded_rect(s11, x_pos, Inches(2.1), card_w, card_h, SOFT_BLUE)
        add_rect(s11, x_pos, Inches(2.1), card_w, Inches(0.12), BRAND_BLUE)
        
        # 텍스트
        add_text(s11, x_pos + Inches(0.2), Inches(2.4), card_w - Inches(0.4), Inches(0.5), title, 18, color=BRAND_BLUE, bold=True, font=FONT_TITLE)
        add_text(s11, x_pos + Inches(0.2), Inches(3.0), card_w - Inches(0.4), Inches(0.5), path, 11, color=GRAY_TEXT, font=FONT_BODY)
        add_text(s11, x_pos + Inches(0.2), Inches(3.6), card_w - Inches(0.4), Inches(2.6), desc, 13, color=DARK_TEXT, font=FONT_BODY)
        
        if idx < 3:
            # 연결 선 데코
            arrow_x = x_pos + card_w + Inches(0.05)
            ln = s11.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, arrow_x, Inches(4.25), arrow_x + Inches(0.2), Inches(4.25))
            ln.line.color.rgb = BRAND_BLUE
            ln.line.width = Pt(1.5)

    add_footer(s11, 11)

    # ----------------------------------------------------
    # SLIDE 12: DB 구성 (DATABASE)
    # ----------------------------------------------------
    s12 = new_slide()
    add_header(s12, "DATABASE SCHEMA", "Prisma & Supabase DB 관계 모델링 (ERD)")
    
    # Left description card
    add_rounded_rect(s12, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5), SOFT_BLUE)
    add_rect(s12, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.12), BRAND_BLUE)
    
    add_text(s12, Inches(1.1), Inches(2.4), Inches(5.0), Inches(0.5), "Prisma Schema & Cascade 제약", 18, color=BRAND_BLUE, bold=True, font=FONT_TITLE)
    
    db_bullets = [
        "• **Supabase PostgreSQL** 클라우드 데이터베이스를 기반으로 고가용성 예약 인프라 구축.",
        "• **User ── Reservation ── ReservationItem**의 정규화된 3대 핵심 테이블 구성.",
        "• 예약 삭제(DELETE) 발생 시, 연결된 ReservationItem(주문요리) 데이터가 Cascade 룰에 의해 안전하고 정밀하게 연쇄 물리 자동 소멸."
    ]
    for idx, bullet in enumerate(db_bullets):
        add_text(s12, Inches(1.1), Inches(3.1) + Inches(0.9) * idx, Inches(5.0), Inches(0.8), bullet, 14, color=DARK_TEXT, font=FONT_BODY)

    # Right ERD Image
    erd_artifact_path = r"C:\Users\yhm2409\.gemini\antigravity\brain\514b8155-baf4-4d6f-923a-5def5c7653b5\media__1784168928988.png"
    if os.path.exists(erd_artifact_path):
        s12.shapes.add_picture(erd_artifact_path, Inches(6.9), Inches(2.0), width=Inches(5.7))

    add_footer(s12, 12)

    # ----------------------------------------------------
    # SLIDE 13: 05. 서비스 이용 시연 간지
    # ----------------------------------------------------
    create_section_slide("05", "서비스 이용 시연", 13)

    # ----------------------------------------------------
    # SLIDE 14: 시연 시나리오 1 (중복 예약 방지) - 3단계 가로 카드 도식화
    # ----------------------------------------------------
    s14 = new_slide()
    add_header(s14, "DEMO SCENARIO 01", "실시간 테이블 선점 락 & 중복 예약 방지")
    
    lock_steps = [
        ("일정 클릭 즉시 조회", "사용자가 방문 일자와 시간대를 토글하는 즉시 API가 호출되어 DB로부터 기완료된 좌석을 정밀 타격 조회합니다."),
        ("SeatMap 실시간 락", "선점된 테이블 ID(Table 1~6)는 화면 상에서 즉시 회색으로 강제 비활성화되어 다른 세션의 터치/선택을 차단합니다."),
        ("트랜잭션 충돌 원천 제어", "예약 완료 POST 전송 시 서버가 중복 여부를 최종 재확인함으로써 다중 동시 선점 문제를 완벽 해결합니다.")
    ]
    
    step_w = Inches(3.6)
    step_h = Inches(3.8)
    s_left = Inches(0.8)
    s_gap = Inches(0.4)
    
    for idx, (title, desc) in enumerate(lock_steps):
        x_pos = s_left + (step_w + s_gap) * idx
        add_rounded_rect(s14, x_pos, Inches(2.4), step_w, step_h, SOFT_BLUE)
        
        # 데코 띠
        add_rect(s14, x_pos + Inches(0.3), Inches(2.7), Inches(1.5), Inches(0.4), BRAND_BLUE)
        add_text(s14, x_pos + Inches(0.3), Inches(2.75), Inches(1.5), Inches(0.4), f"STAGE 0{idx+1}", 13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_BODY)
        
        # 텍스트
        add_text(s14, x_pos + Inches(0.3), Inches(3.4), step_w - Inches(0.6), Inches(0.6), title, 17, color=BRAND_BLUE, bold=True, font=FONT_TITLE)
        add_text(s14, x_pos + Inches(0.3), Inches(4.1), step_w - Inches(0.6), Inches(1.8), desc, 14, color=DARK_TEXT, font=FONT_BODY)
        
        if idx < 2:
            arrow_x = x_pos + step_w + Inches(0.05)
            ln = s14.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, arrow_x, Inches(4.3), arrow_x + Inches(0.3), Inches(4.3))
            ln.line.color.rgb = BRAND_BLUE
            ln.line.width = Pt(2.0)

    add_footer(s14, 14)

    # ----------------------------------------------------
    # SLIDE 15: 시연 시나리오 2 (AI 챗봇 및 어드민 연동) - 2분할 카드 도식화
    # ----------------------------------------------------
    s15 = new_slide()
    add_header(s15, "DEMO SCENARIO 02", "AI 상담 챗봇 및 어드민 실시간 삭제 동기화")
    
    # Left Card: AI Chatbot
    add_rounded_rect(s15, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5), SOFT_BLUE)
    add_rect(s15, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.12), BRAND_BLUE)
    add_text(s15, Inches(1.1), Inches(2.4), Inches(5.0), Inches(0.5), "매장 데이터 동기화 AI 챗봇", 20, color=BRAND_BLUE, bold=True, font=FONT_TITLE)
    
    chat_points = [
        "• 실제 예매 화면의 한우 채끝 코스 요리 메뉴, 6개 테이블의 특징, 무료 발레파킹 정보를 Llama-3.3 LLM에 100% 바인딩.",
        "• 소비자가 '발레파킹 요금이 얼마인가요?', '조용한 룸 예약 가능한가요?' 질문 시 매장 실제 사양에 완전히 밀착된 정밀 응대 시연."
    ]
    for idx, pt in enumerate(chat_points):
        add_text(s15, Inches(1.1), Inches(3.1) + Inches(1.5) * idx, Inches(5.0), Inches(1.4), pt, 14, color=DARK_TEXT, font=FONT_BODY)

    # Right Card: Admin Link
    add_rounded_rect(s15, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.5), SOFT_BLUE)
    add_rect(s15, Inches(6.9), Inches(2.0), Inches(5.6), Inches(0.12), BRAND_BLUE)
    add_text(s15, Inches(7.2), Inches(2.4), Inches(5.0), Inches(0.5), "어드민 실시간 삭제 & 좌석 복원", 20, color=BRAND_BLUE, bold=True, font=FONT_TITLE)
    
    admin_points = [
        "• 사장님이 관리자 탭에서 승인된 예약을 물리 영구 삭제하는 즉시 유저 마이페이지 내 예약 현황 리스트에서 삭제 연계.",
        "• 삭제 처리 즉시 해당 예약 슬롯의 좌석 락이 즉각 해제되어 SeatMap 상에서 예약 가능한 녹색 상태로 동적 복구 시연."
    ]
    for idx, pt in enumerate(admin_points):
        add_text(s15, Inches(7.2), Inches(3.1) + Inches(1.5) * idx, Inches(5.0), Inches(1.4), pt, 14, color=DARK_TEXT, font=FONT_BODY)

    add_footer(s15, 15)

    # ----------------------------------------------------
    # SLIDE 16: Q&A / 감사합니다 (마무리 표지)
    # ----------------------------------------------------
    s16 = new_slide()
    # 대형 블루 카드
    add_rounded_rect(s16, Inches(0.8), Inches(1.5), Inches(11.73), Inches(4.2), fill_color=BRAND_BLUE)
    
    add_text(s16, Inches(1.4), Inches(2.5), Inches(10.5), Inches(1.2), "경청해 주셔서 감사합니다", 46, color=WHITE, bold=True, font=FONT_TITLE)
    add_text(s16, Inches(1.4), Inches(3.8), Inches(10.5), Inches(0.8), "라롬 (L'AROME) 식당 예약 플랫폼 Q&A", 20, color=SOFT_BLUE, bold=True, font=FONT_BODY)
    
    add_footer(s16, 16)

    # 저장 경로 (최종 화이트 & 블루 스타일)
    dest_path = r"c:\Users\yhm2409\Desktop\map\L_AROME_Presentation_BlueStyle.pptx"
    prs.save(dest_path)
    
    # 임시 변환 이미지 정리
    if os.path.exists(temp_png):
        os.remove(temp_png)
    print("White & Blue Style PPT built successfully!")

if __name__ == "__main__":
    build()
