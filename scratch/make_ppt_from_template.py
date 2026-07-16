import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

def build_presentation():
    template_path = "나만의서비스_발표자료_템플릿.pptx"
    if not os.path.exists(template_path):
        print("Template file not found!")
        return
        
    prs = Presentation(template_path)
    
    # 폰트 컬러 설정용
    dark_gray = RGBColor(78, 89, 104)
    black = RGBColor(25, 31, 40)
    toss_blue = RGBColor(49, 130, 246)
    
    # 헬퍼 함수: 텍스트 상자 새로 추가하여 내용 기입
    def add_bullet_points(slide, points, left, top, width, height, font_size=15):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        for idx, pt in enumerate(points):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = pt
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(font_size)
            if pt.startswith("•") or pt.startswith("  -") or pt.startswith("   -"):
                p.font.bold = pt.startswith("•")
                p.font.color.rgb = black if pt.startswith("•") else dark_gray
            else:
                p.font.bold = True
                p.font.color.rgb = black
            p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 1: 표지
    # ----------------------------------------------------
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.has_text_frame:
            txt = shape.text
            if "MY SERVICE" in txt:
                shape.text = "L'AROME FINE DINING"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = toss_blue
            elif " " in txt:
                shape.text = "라롬 (L'AROME) 식당 예약 및 AI 상담 플랫폼"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black
            elif "̸" in txt or "ǥ" in txt:
                shape.text = "발표자: [이름]  |  개발 프로젝트 완료 보고"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.color.rgb = dark_gray

    # ----------------------------------------------------
    # SLIDE 2: CONTENTS (목차)
    # ----------------------------------------------------
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if shape.has_text_frame:
            txt = shape.text
            if " Ұ" in txt:
                shape.text = "서비스 소개"
            elif "  " in txt:
                shape.text = "서비스 구현 범위"
            elif " UI  ۷" in txt:
                shape.text = "서비스 UI 및 레퍼런스"
            elif "鿣  DB " in txt:
                shape.text = "백엔드 및 DB 구성 (ERD)"
            elif " ̿ ÿ" in txt:
                shape.text = "서비스 이용 시연"
            
            for p in shape.text_frame.paragraphs:
                p.font.name = "Malgun Gothic"
                p.font.bold = True

    # ----------------------------------------------------
    # SLIDE 3: 간지 1 (서비스 소개)
    # ----------------------------------------------------
    s3 = prs.slides[2]
    for shape in s3.shapes:
        if shape.has_text_frame:
            if " Ұ" in shape.text:
                shape.text = "서비스 소개"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    # ----------------------------------------------------
    # SLIDE 4: SERVICE INTRO (기획 배경)
    # ----------------------------------------------------
    s4 = prs.slides[3]
    for shape in s4.shapes:
        if shape.has_text_frame:
            if "ȹ" in shape.text:
                shape.text = "기획 배경 및 해결하고자 하는 문제"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black
                    
    # 상세 내용 추가
    s4_points = [
        "• 바쁜 현대인의 예약 대기 리소스 최적화",
        "  - 식당 입장을 위해 장시간 오프라인 대기하는 시간과 불편을 근본적으로 극복합니다.",
        "\n• 파인다이닝 예약 부도(No-Show) 방지",
        "  - 사전에 방문 일정과 세부 식사 메뉴를 조율/선택 완료하여 예약 신뢰성을 향상시킵니다.",
        "\n• 정보 불일치 해소",
        "  - 창가석, 룸, 테라스 등 좌석의 실제 형태를 미리 확인하지 못해 발생하던 현장 불만족을 방지합니다."
    ]
    add_bullet_points(s4, s4_points, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5), font_size=16)

    # ----------------------------------------------------
    # SLIDE 5: SERVICE INTRO (서비스 가치)
    # ----------------------------------------------------
    s5 = prs.slides[4]
    for shape in s5.shapes:
        if shape.has_text_frame:
            if "ٽ" in shape.text:
                shape.text = "서비스의 가치 및 차별점"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    s5_points = [
        "• 원스톱 미식 예약 여정",
        "  - 인원/일정 선택 ➡️ 2D 지도 기반 좌석 테이블 선점 ➡️ 사전 메뉴 주문의 매끄러운 3단계 동선 설계.",
        "\n• 직관적 2D 좌석 배치도 (SeatMap)",
        "  - 창가석, 프라이빗 단체 룸, 테라스 뷰 등 좌석 성격을 시각적으로 인지하고 원하는 자리를 미리 확보합니다.",
        "\n• 식당 데이터 100% 동기화 AI 컨시어지",
        "  - 24시간 실시간 소비자 챗봇 상담을 통해 매장의 전화 문의 리소스를 획기적으로 낮춥니다."
    ]
    add_bullet_points(s5, s5_points, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5), font_size=16)

    # ----------------------------------------------------
    # SLIDE 6: 간지 2 (서비스 구현 범위)
    # ----------------------------------------------------
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if shape.has_text_frame:
            if "  " in shape.text:
                shape.text = "서비스 구현 범위"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    # ----------------------------------------------------
    # SLIDE 7: SCOPE (구현 범위)
    # ----------------------------------------------------
    s7 = prs.slides[6]
    for shape in s7.shapes:
        if shape.has_text_frame:
            if " " in shape.text:
                shape.text = "Client, Server, AI 구현 명세"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    s7_points = [
        "• 클라이언트 (Client UI/UX Scope)",
        "  - 모바일 반응형 뷰포트 기기 프레임 및 하단 Sticky 네비게이션 탭바 탑재.",
        "  - 3단계 예약 진행 동선 및 장바구니 총 금액 실시간 정산 기능 구현.",
        "\n• 백엔드 서버 (Backend API Scope)",
        "  - Next.js API Routes 기반 유저 회원가입/로그인 및 인증 관리.",
        "  - 좌석의 실시간 중복 차단을 보장하는 날짜/시간대별 점유 조회 API 구현.",
        "  - 사장님 관리 전용 데스크톱 가로형 어드민 뷰포트 연계 및 예약 승인/거절/영구 삭제 API.",
        "\n• 인공지능 (AI Scope)",
        "  - GROQ SDK 및 Llama-3.3 LLM 연계 실시간 매장 자동화 상담 챗봇 탑재."
    ]
    add_bullet_points(s7, s7_points, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5), font_size=14)

    # ----------------------------------------------------
    # SLIDE 8: 간지 3 (서비스 UI 및 레퍼런스)
    # ----------------------------------------------------
    s8 = prs.slides[7]
    for shape in s8.shapes:
        if shape.has_text_frame:
            if " UI" in shape.text:
                shape.text = "서비스 UI 및 레퍼런스"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    # ----------------------------------------------------
    # SLIDE 9: USER FLOW (작업 프로세스)
    # ----------------------------------------------------
    s9 = prs.slides[8]
    for shape in s9.shapes:
        if shape.has_text_frame:
            if " ÷ο" in shape.text:
                shape.text = "작업 프로세스 및 예약 사용자 흐름"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    s9_points = [
        "• STEP 01. 방문 일정 설정",
        "  - 모바일 캘린더에서 예약일자, 인원수, 디너 타임 시간대를 정밀 조율합니다.",
        "\n• STEP 02. 가상 테이블 선택",
        "  - 선택 일정의 점유 락을 체크하여 이미 찬 테이블을 제외한 나머지 좌석(창가석, 룸 등)을 지정합니다.",
        "\n• STEP 03. 요리 코스 주문",
        "  - 매칭 코스 요리 메뉴판에서 원하는 세트를 장바구니에 담고 모의 결제 완료를 진행합니다.",
        "\n• STEP 04. 내 예약 확인 및 취소",
        "  - 마이페이지 내역에서 승인대기/완료 상태를 모니터링하고 원클릭 삭제 연계를 진행합니다."
    ]
    add_bullet_points(s9, s9_points, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5), font_size=14)

    # ----------------------------------------------------
    # SLIDE 10: REFERENCE (UI 레퍼런스)
    # ----------------------------------------------------
    s10 = prs.slides[9]
    for shape in s10.shapes:
        if shape.has_text_frame:
            if " ۷" in shape.text:
                shape.text = "UI 레퍼런스 및 토스 스타일 가이드"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    s10_points = [
        "• 토스(TOSS) UI 디자인 스타일 상속",
        "  - 군더더기 없는 플랫 둥근 카드 레이아웃과 뚜렷한 Toss Blue 테마 채용.",
        "  - 모바일 조작 안정성을 보장하는 엄지 영역(Bottom Nav Bar) 네비게이션 가이드.",
        "\n• 애자일 점진적 작업 프로세스",
        "  - 뼈대 구현 후 도메인 확장 피드백 수렴.",
        "  - 식당 ➡️ 영화 ➡️ 식당 컨셉 롤백에 따른 챗봇 지식 베이스 복원 및 UI 통합 정교화."
    ]
    add_bullet_points(s10, s10_points, Inches(0.8), Inches(2.2), Inches(7.5), Inches(4.5), font_size=15)

    # WebP ➡️ PNG 변환 및 삽입
    try:
        toss_webp = r"c:\Users\yhm2409\Desktop\map\toss\imgi_61_73ab268a060611ed89c61329738d4db8_small.webp"
        temp_png = r"c:\Users\yhm2409\Desktop\map\toss\temp_toss_ref_temp.png"
        if os.path.exists(toss_webp):
            with Image.open(toss_webp) as img:
                img.save(temp_png, "PNG")
            s10.shapes.add_picture(temp_png, Inches(8.8), Inches(1.8), width=Inches(3.8))
    except Exception as e:
        print("Toss Image insertion failed:", e)

    # ----------------------------------------------------
    # SLIDE 11: 간지 4 (백엔드 및 DB 구성)
    # ----------------------------------------------------
    s11 = prs.slides[10]
    for shape in s11.shapes:
        if shape.has_text_frame:
            if "鿣" in shape.text:
                shape.text = "백엔드 및 DB 구성"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    # ----------------------------------------------------
    # SLIDE 12: DATABASE (ERD)
    # ----------------------------------------------------
    s12 = prs.slides[11]
    for shape in s12.shapes:
        if shape.has_text_frame:
            if "ERD" in shape.text:
                shape.text = "데이터베이스 관계형 스키마 설계 (ERD)"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    s12_points = [
        "• User ── Reservation ── ReservationItem",
        "  - 3대 핵심 모델 간의 깔끔한 1:N 릴레이션 관계.",
        "\n• 데이터 참조 무결성 확보",
        "  - Prisma Schema 파일 내에 Cascade Delete 외래키 규칙 명시.",
        "  - 예약 취소/물리 삭제 시, 해당 예약에 속했던 하위 주문 요리 품목들이 연쇄 삭제되어 쓰레기 데이터 누적 방지."
    ]
    add_bullet_points(s12, s12_points, Inches(0.8), Inches(2.2), Inches(6.0), Inches(4.5), font_size=15)

    # ERD 이미지 삽입
    erd_artifact_path = r"C:\Users\yhm2409\.gemini\antigravity\brain\514b8155-baf4-4d6f-923a-5def5c7653b5\media__1784168928988.png"
    if os.path.exists(erd_artifact_path):
        s12.shapes.add_picture(erd_artifact_path, Inches(7.0), Inches(2.0), width=Inches(5.7))

    # ----------------------------------------------------
    # SLIDE 13: TECH STACK (기술 스택)
    # ----------------------------------------------------
    s13 = prs.slides[12]
    for shape in s13.shapes:
        if shape.has_text_frame:
            if " " in shape.text:
                shape.text = "기술 스택 및 인프라 아키텍처"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    s13_points = [
        "• 프론트엔드 (Frontend)",
        "  - React 19, TypeScript, Vanilla CSS (리치 미세 애니메이션 및 컴포넌트 튜닝)",
        "\n• 백엔드 및 인프라 (Backend & Infrastructure)",
        "  - Next.js 16 (App Router & Serverless Route Handlers API)",
        "  - Prisma ORM & Supabase PostgreSQL (물리 원격 클라우드 DB 인프라)",
        "\n• 인공지능 API (AI Inference API)",
        "  - GROQ SDK & Llama-3.3 70B 모델 결합 (매장 동기화 컨텍스트 AI 챗봇)"
    ]
    add_bullet_points(s13, s13_points, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5), font_size=15)

    # ----------------------------------------------------
    # SLIDE 14: 간지 5 (서비스 이용 시연)
    # ----------------------------------------------------
    s14 = prs.slides[13]
    for shape in s14.shapes:
        if shape.has_text_frame:
            if " ̿" in shape.text:
                shape.text = "서비스 이용 시연"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    # ----------------------------------------------------
    # SLIDE 15: DEMO (시연 시나리오)
    # ----------------------------------------------------
    s15 = prs.slides[14]
    for shape in s15.shapes:
        if shape.has_text_frame:
            if "̿ ÿ" in shape.text:
                shape.text = "핵심 가치 시연 시나리오"
                for p in shape.text_frame.paragraphs:
                    p.font.name = "Malgun Gothic"
                    p.font.bold = True
                    p.font.color.rgb = black

    s15_points = [
        "• 시나리오 1. 실시간 중복 예약 방지 락 (Seat Lock)",
        "  - 특정 날짜와 시간대에 테이블(Table 1)을 예약 접수하는 즉시, 다른 세션에서 해당 좌석이 SeatMap 지도 상에서 회색 처리(비활성화)되어 중복 예약 충돌이 원천 방지되는 흐름 시연.",
        "\n• 시나리오 2. 100% 매장 동기화 AI 챗봇",
        "  - 메뉴 명세, 6개 테이블 상세, 주차 가이드 등에 대해 실시간 자연어 질문 시, 백엔드 컨텍스트 연동으로 오류 없는 응대 시연.",
        "\n• 시나리오 3. 사장님(어드민) 모드 실시간 삭제 & 좌석 복원",
        "  - 어드민 뷰에서 수동 승인/거절 및 삭제 처리 시, 즉시 SeatMap에서 테이블 잠금이 풀려 신규 예약이 가능해지는 실시간 연계 시연."
    ]
    add_bullet_points(s15, s15_points, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5), font_size=14)

    # Save to L_AROME_Presentation_Final.pptx
    dest_path = r"c:\Users\yhm2409\Desktop\map\L_AROME_Presentation_Final.pptx"
    prs.save(dest_path)
    
    # clean up temp image
    if os.path.exists(temp_png):
        os.remove(temp_png)
    print("Template PPT built successfully!")

if __name__ == "__main__":
    build_presentation()
