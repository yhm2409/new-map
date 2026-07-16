# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from PIL import Image

# 템플릿의 색상 조합
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x13, 0x4E, 0x4A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
LIGHT_TEAL_LINE = RGBColor(0xB8, 0xD3, 0xD0)
FONT = "Malgun Gothic"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

def new_slide():
    return prs.slides.add_slide(blank)

def set_bg_white(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

def add_textbox(slide, left, top, width, height, text, size, color=TEAL,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return box

def add_bullet_points(slide, points, left, top, width, height, font_size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    for idx, pt in enumerate(points):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = pt
        p.font.name = FONT
        p.font.size = Pt(font_size)
        if pt.startswith("•") or pt.startswith("  -") or pt.startswith("   -"):
            p.font.bold = pt.startswith("•")
            p.font.color.rgb = TEAL if pt.startswith("•") else dark_gray_fallback
        else:
            p.font.bold = True
            p.font.color.rgb = TEAL
        p.space_after = Pt(8)

dark_gray_fallback = RGBColor(78, 89, 104)

def add_line(slide, left, top, width, color=LIGHT_TEAL_LINE, weight=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln

def add_rect(slide, left, top, width, height, fill_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    if fill_color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def add_page_footer(slide, num):
    add_textbox(slide, SW - Inches(1.2), SH - Inches(0.6), Inches(0.9), Inches(0.4),
                str(num), 11, color=ORANGE, align=PP_ALIGN.RIGHT)

def add_header(slide, kicker, title):
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4), kicker, 13, color=ORANGE, bold=True)
    add_textbox(slide, Inches(0.68), Inches(0.85), Inches(11.5), Inches(0.9), title, 30, color=TEAL, bold=True)
    add_line(slide, Inches(0.7), Inches(1.7), Inches(11.9), color=LIGHT_TEAL_LINE, weight=1.25)

def content_slide(kicker, title, num):
    s = new_slide()
    set_bg_white(s)
    add_header(s, kicker, title)
    add_page_footer(s, num)
    return s

def section_slide(no, title, num):
    s = new_slide()
    set_bg_white(s)
    add_textbox(s, Inches(0.9), Inches(3.2), Inches(2), Inches(0.6), no, 20, color=ORANGE, bold=True)
    add_textbox(s, Inches(0.85), Inches(3.6), Inches(11.5), Inches(1), title, 38, color=TEAL, bold=True)
    add_line(s, Inches(0.9), Inches(4.6), Inches(2.4), color=ORANGE, weight=1.5)
    add_page_footer(s, num)
    return s

def build():
    # ---------- 1. 표지 ----------
    s1 = new_slide()
    set_bg_white(s1)
    add_rect(s1, 0, SH - Inches(0.12), SW, Inches(0.12), fill_color=TEAL)
    add_textbox(s1, Inches(0.9), Inches(2.6), Inches(10), Inches(0.5), "L'AROME FINE DINING", 16, color=ORANGE, bold=True)
    add_textbox(s1, Inches(0.85), Inches(3.0), Inches(11.5), Inches(1.4), "라롬 식당 예약 및 AI 상담 플랫폼", 48, color=TEAL, bold=True)
    add_line(s1, Inches(0.9), Inches(4.35), Inches(3.2), color=ORANGE, weight=1.5)
    add_textbox(s1, Inches(0.9), Inches(4.55), Inches(8), Inches(0.5), "발표자: [이름]  ·  개발 프로젝트 완료 보고", 16, color=TEAL)
    add_page_footer(s1, 1)

    # ---------- 2. 목차 ----------
    s2 = content_slide("CONTENTS", "목차", 2)
    items = [
        ("01", "서비스 소개"),
        ("02", "서비스 구현 범위"),
        ("03", "서비스 UI 및 레퍼런스 (작업 프로세스)"),
        ("04", "백엔드 및 DB 구성 (ERD)"),
        ("05", "서비스 이용 시연")
    ]
    top = Inches(2.1)
    for idx, (no, label) in enumerate(items):
        y = top + Inches(0.95) * idx
        add_textbox(s2, Inches(0.8), y, Inches(1.2), Inches(0.7), no, 24, color=ORANGE, bold=True)
        add_textbox(s2, Inches(1.9), y + Inches(0.04), Inches(9), Inches(0.6), label, 20, color=TEAL, bold=True)
        add_line(s2, Inches(0.8), y + Inches(0.75), Inches(10.8), color=LIGHT_TEAL_LINE, weight=1.0)

    # ---------- 3. 1. 서비스 소개 간지 ----------
    section_slide("01", "서비스 소개", 3)

    # ---------- 4. 서비스 소개 - 기획 배경 ----------
    s4 = content_slide("SERVICE INTRO", "기획 배경 및 해결하고자 하는 문제", 4)
    s4_points = [
        "• 현대인들의 예약 대기 피로도 극복",
        "  - 식당 입장을 위해 오랜 시간 현장에서 무작정 대기하는 비효율성을 해소합니다.",
        "\n• 파인다이닝의 예약 부도(No-Show) 차단",
        "  - 일정 확정과 세부 코스 메뉴 사전 선결제/선택을 유도하여 노쇼 리스크를 제로화합니다.",
        "\n• 테이블 불확실성 해소",
        "  - 고객이 선호하는 창가석, 룸 등의 실제 공간 위치를 지도 상에서 인지하여 현장 갈등을 극복합니다."
    ]
    add_bullet_points(s4, s4_points, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5), font_size=15)

    # ---------- 5. 서비스 소개 - 서비스 가치 ----------
    s5 = content_slide("SERVICE VALUE", "서비스 가치 및 제공 기능 개요", 5)
    s5_points = [
        "• 원스톱 미식 예약 여정",
        "  - 날짜/시간 지정 ➡️ 2D 좌석 지도 지정 ➡️ 코스 요리 메뉴 선택(장바구니) 및 모의 결제 연동.",
        "\n• 시각적 2D 좌석 배치도 (SeatMap)",
        "  - 매장의 창가석, 프라이빗 룸, 테라스 등 공간 배치를 입체적으로 확인하고 선점 가능한 환경 제공.",
        "\n• AI 컨시어지 상담 연계",
        "  - 매장 지식 데이터베이스와 유기적으로 소통하여 24시간 정확한 소비자 안내 제공."
    ]
    add_bullet_points(s5, s5_points, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5), font_size=15)

    # ---------- 6. 2. 서비스 구현 범위 간지 ----------
    section_slide("02", "서비스 구현 범위", 6)

    # ---------- 7. 구현 범위 - SCOPE ----------
    s7 = content_slide("DEVELOPMENT SCOPE", "Client, Server, AI 구현 명세", 7)
    s7_points = [
        "• 클라이언트 (Client UI/UX)",
        "  - 모바일 반응형 뷰포트 기기 프레임 및 하단 Sticky 네비게이션 탭바 탑재.",
        "  - 3단계 예약 플로우 및 메뉴판 고정 스크롤 박스, 장바구니 총 금액 실시간 합산.",
        "\n• 백엔드 서버 (Backend API & Admin)",
        "  - Next.js API Routes 기반 유저 회원가입/로그인 및 인증 관리 API.",
        "  - 날짜/시간대별 실시간 예약 테이블 점유 조회 및 중복 선점 차단 락 API.",
        "  - 관리자 전용 데스크톱 가로 뷰포트 대시보드 및 예약 승인/거절/영구 삭제 API.",
        "\n• 인공지능 (AI Scope)",
        "  - GROQ SDK 및 Llama-3.3 70B LLM 서버 연동 매장 자동화 상담 챗봇."
    ]
    add_bullet_points(s7, s7_points, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.8), font_size=14)

    # ---------- 8. 3. 서비스 UI 및 레퍼런스 간지 ----------
    section_slide("03", "서비스 UI 및 레퍼런스 (작업 프로세스)", 8)

    # ---------- 9. UI 및 레퍼런스 - REFERENCE ----------
    s9 = content_slide("UI & REFERENCE", "TOSS 디자인 스타일 및 애자일 작업 프로세스", 9)
    s9_points = [
        "• 토스(TOSS) UI 디자인 스타일 계승",
        "  - 군더더기 없는 플랫 둥근 카드 레이아웃과 명료한 Toss Blue 테마 적용.",
        "  - 모바일 환경에서의 엄지 영역 접근성(하단 네비게이션 가이드) 극대화.",
        "\n• 기민한(Agile) 점진적 작업 프로세스",
        "  - DB 관계 모델링 ➡️ 핵심 예약 모듈 빌드 ➡️ AI 챗봇 연동 및 고도화.",
        "  - 영화 예매 도메인 피드백 수렴 후 식당 예약 뼈대로 신속히 원상 복구 및 컨텍스트 동기화 수행."
    ]
    add_bullet_points(s9, s9_points, Inches(0.8), Inches(2.0), Inches(7.5), Inches(4.5), font_size=14)

    # Toss WebP ➡️ PNG 변환 및 삽입
    try:
        toss_webp = r"c:\Users\yhm2409\Desktop\map\toss\imgi_61_73ab268a060611ed89c61329738d4db8_small.webp"
        temp_png = r"c:\Users\yhm2409\Desktop\map\toss\temp_toss_ref_teal.png"
        if os.path.exists(toss_webp):
            with Image.open(toss_webp) as img:
                img.save(temp_png, "PNG")
            s9.shapes.add_picture(temp_png, Inches(8.8), Inches(2.0), width=Inches(3.8))
    except Exception as e:
        print("Toss Image insertion failed:", e)

    # ---------- 10. 4. 백엔드 및 DB 구성 간지 ----------
    section_slide("04", "백엔드 및 DB 구성", 10)

    # ---------- 11. 백엔드 구성 - BACKEND ----------
    s11 = content_slide("BACKEND ARCHITECTURE", "Next.js API Routes 및 AI 추론 연동", 11)
    s11_points = [
        "• RESTful API API Routes 설계",
        "  - `/api/auth/login` & `/api/auth/signup`: 인증 및 세션 토큰 제어.",
        "  - `/api/reservations`: 날짜/시간대 선택에 따른 좌석 실시간 예약 상태 체크 및 POST 접수.",
        "  - `/api/reservations/[id]`: 어드민 모드 예약 승인/거절 상태 변경 및 완전 Cascade 삭제.",
        "\n• AI 추론 API (`/api/chat`) 연계",
        "  - 브라우저 클라이언트가 아닌 Node.js 백엔드 환경에서 GROQ API와 통신하여 API Key 유출 방지.",
        "  - 시스템 프롬프트(식당 정보 100% 동기화)를 내부 바인딩하여 안전성 극대화."
    ]
    add_bullet_points(s11, s11_points, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5), font_size=15)

    # ---------- 12. DB 구성 - DATABASE ----------
    s12 = content_slide("DATABASE SCHEMA", "Prisma & Supabase DB 관계 모델링 (ERD)", 12)
    s12_points = [
        "• Prisma ORM & Supabase Postgres 결합",
        "  - 데이터 신뢰성과 관계 무결성을 확보한 PostgreSQL 클라우드 DB 연동.",
        "\n• 외래키 관계 및 연쇄 삭제 (Cascade Delete)",
        "  - User, Reservation, ReservationItem 간의 1:N 구조 설정.",
        "  - 부모 예약 레코드 삭제 시 자식 주문 요리 품목들이 연쇄 삭제되어 쓰레기 데이터 방지."
    ]
    add_bullet_points(s12, s12_points, Inches(0.8), Inches(2.0), Inches(6.0), Inches(4.5), font_size=15)

    # ERD 이미지 삽입
    erd_artifact_path = r"C:\Users\yhm2409\.gemini\antigravity\brain\514b8155-baf4-4d6f-923a-5def5c7653b5\media__1784168928988.png"
    if os.path.exists(erd_artifact_path):
        s12.shapes.add_picture(erd_artifact_path, Inches(7.0), Inches(2.0), width=Inches(5.7))

    # ---------- 13. 5. 서비스 이용 시연 간지 ----------
    section_slide("05", "서비스 이용 시연", 13)

    # ---------- 14. 서비스 이용 시연 - 중복 예약 방지 ----------
    s14 = content_slide("DEMO SCENARIO 01", "실시간 테이블 선점 락 & 중복 예약 방지", 14)
    s14_points = [
        "• 일정/시간 선택 즉시 예약 상태 체크",
        "  - 사용자가 디너 타임 일정을 결정하는 순간 서버 API가 동적으로 작동하여 점유 테이블 ID를 리턴합니다.",
        "\n• SeatMap 좌석 비활성화 동적 락",
        "  - 이미 예약된 테이블(Table 1~6)은 화면 상에서 회색 처리(비활성화)되어 마우스/터치 선택이 불가능해집니다.",
        "\n• DB 트랜잭션 충돌 차단",
        "  - 다중 사용자의 동시 좌석 선점 오류를 원천 차단하여 안전한 예약 접수를 처리합니다."
    ]
    add_bullet_points(s14, s14_points, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5), font_size=15)

    # ---------- 15. 서비스 이용 시연 - AI 챗봇 및 어드민 ----------
    s15 = content_slide("DEMO SCENARIO 02", "AI 상담 챗봇 & 어드민 삭제 실시간 복원", 15)
    s15_points = [
        "• 매장 스펙 100% 동기화 AI 컨시어지 챗봇",
        "  - 실제 코스 메뉴(스테이크 등), 좌석 성격, 주차 정보 등에 대해 GROQ Llama-3.3 LLM을 통해 정교한 응대 제공.",
        "\n• 어드민 예약 처리 연동 실시간 좌석 복원",
        "  - 사장님이 어드민 모드에서 승인 예약 취소(Cascade 완전 물리 삭제)를 누르는 순간 DB가 클리어링됩니다.",
        "  - 삭제와 즉시 SeatMap 상에서 해당 좌석 비활성화 락이 즉시 해제되어 다시 예약 가능 상태로 복원됩니다."
    ]
    add_bullet_points(s15, s15_points, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5), font_size=14)

    # ---------- 16. 감사합니다 (Q&A) ----------
    s16 = new_slide()
    set_bg_white(s16)
    add_textbox(s16, Inches(1.0), Inches(2.5), Inches(11.3), Inches(3.0), "경청해 주셔서 감사합니다", 44, color=TEAL, bold=True)
    add_textbox(s16, Inches(1.0), Inches(3.6), Inches(11.3), Inches(1.0), "L'AROME Fine Dining 예약 플랫폼 Q&A", 20, color=ORANGE, bold=True)
    add_page_footer(s16, 16)

    # 저장 경로
    dest_path = r"c:\Users\yhm2409\Desktop\map\L_AROME_Presentation_TealStyle.pptx"
    prs.save(dest_path)
    
    # 임시 이미지 정리
    if os.path.exists(temp_png):
        os.remove(temp_png)
    print("Teal & Orange Style PPT built successfully!")

if __name__ == "__main__":
    build()
